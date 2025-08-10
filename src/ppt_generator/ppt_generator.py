import os
import sys

sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from utils.file_utils import ensure_output_dir

def set_font_size_12pt(text_frame):
    """
    Sets the font size to 12pt for all paragraphs in a text frame.
    
    Args:
        text_frame: The text frame to modify
    """
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(12)

def create_ppt_for_story(story_path, images_dir, output_dir, max_words_per_slide=450):
    """
    Creates a PPTX file for a single story.

    Args:
        story_path (str): Path to the story text file.
        images_dir (str): Directory containing images for the story.
        output_dir (str): Directory to save the generated PPTX.
        max_words_per_slide (int): Maximum words per slide.
    """
    ensure_output_dir(output_dir)
    story_name = os.path.splitext(os.path.basename(story_path))[0]
    ppt_filename = f"{story_name}.pptx"
    ppt_path = os.path.join(output_dir, ppt_filename)

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    # Add title and subtitle
    title.text = story_name.replace("_", " ").title()
    with open(story_path, 'r', encoding='utf-8') as file:
        story = file.read()
    subtitle.text = story[:100]  # Short description
    
    # Set font size to 12pt for title and subtitle
    set_font_size_12pt(title.text_frame)
    set_font_size_12pt(subtitle.text_frame)

    # Add content slides
    words = story.split()
    for i in range(0, len(words), max_words_per_slide):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        content = slide.shapes.placeholders[1]
        content.text = " ".join(words[i:i + max_words_per_slide])
        content.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        # Set font size to 12pt for content
        set_font_size_12pt(content.text_frame)

    # Add image slide
    image_prefix = os.path.splitext(story_name)[0]
    image_filename = f"{image_prefix}_image_1.png"
    image_path = os.path.join(images_dir, image_filename)
    if os.path.exists(image_path):
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        left = top = Inches(1)
        pic = slide.shapes.add_picture(image_path, left, top, width=Inches(6), height=Inches(4.5))

    # Save PPTX
    prs.save(ppt_path)
    print(f"PPTX created: {ppt_path}")


def combine_ppts(input_dir, output_file):
    """
    Combines multiple PPTX files into a single PPTX file.

    Args:
        input_dir (str): Directory containing individual PPTX files.
        output_file (str): Path to save the combined PPTX file.
    """
    from pptx import Presentation

    prs = Presentation()
    for ppt_file in os.listdir(input_dir):
        if ppt_file.endswith(".pptx"):
            src_prs = Presentation(os.path.join(input_dir, ppt_file))
            for slide in src_prs.slides:
                slide_layout = prs.slide_layouts[0]
                new_slide = prs.slides.add_slide(slide_layout)
                
                # Track text content for the slide
                slide_text = ""
                
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            slide_text += paragraph.text + "\n"
                    elif shape.shape_type == 13:  # Picture
                        # Save image to a temp file and re-insert
                        image = shape.image
                        image_bytes = image.blob
                        import tempfile
                        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
                            tmp_img.write(image_bytes)
                            tmp_img_path = tmp_img.name
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height
                        new_slide.shapes.add_picture(tmp_img_path, left, top, width, height)
                        os.remove(tmp_img_path)
                
                # Add text content to the new slide and set font size to 12pt
                if slide_text.strip():
                    new_slide.shapes[0].text = slide_text
                    set_font_size_12pt(new_slide.shapes[0].text_frame)
                    
    prs.save(output_file)
    print(f"Combined PPTX saved to: {output_file}")

if __name__ == "__main__":
    pass
    # # Test directories (adjust as needed)
    # stories_dir = "data/output/stories"
    # images_dir = "data/output/images"
    # pptx_dir = "data/output/pptx"
    # combined_pptx = "data/output/pptx/combined_stories.pptx"

    # # Ensure output directory exists
    # if not os.path.exists(pptx_dir):
    #     os.makedirs(pptx_dir)

    # # Test: Create PPTX for each story
    # for story_file in os.listdir(stories_dir):
    #     if story_file.endswith(".txt"):
    #         story_path = os.path.join(stories_dir, story_file)
    #         create_ppt_for_story(story_path, images_dir, pptx_dir)

    # # Test: Combine all PPTX files into one
    # combine_ppts(pptx_dir, combined_pptx)
    # print("Testing complete. Check the pptx directory and combined_stories.pptx.")